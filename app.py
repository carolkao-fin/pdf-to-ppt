import io
import zipfile
import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ── helpers ───────────────────────────────────────────────────────────────────

def _emu(pt_val: float) -> int:
    """PDF points → PowerPoint EMUs  (1 pt = 1/72 inch = 12700 EMU)"""
    return int(pt_val * 12700)


def pdf_page_count(pdf_bytes: bytes) -> int:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    n = len(doc)
    doc.close()
    return n


def render_previews(pdf_bytes: bytes, max_pages: int = 6, scale: float = 0.7):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []
    mat = fitz.Matrix(scale, scale)
    for i in range(min(len(doc), max_pages)):
        pix = doc[i].get_pixmap(matrix=mat)
        images.append(pix.tobytes("png"))
    doc.close()
    return images


# ── conversion：PPTX ──────────────────────────────────────────────────────────

def to_pptx_image(pdf_bytes: bytes) -> bytes:
    """每頁轉為高解析度圖片投影片，版面完整保留。"""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    prs = Presentation()
    mat = fitz.Matrix(2, 2)

    for page in doc:
        w_emu = _emu(page.rect.width)
        h_emu = _emu(page.rect.height)
        prs.slide_width = w_emu
        prs.slide_height = h_emu
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pix = page.get_pixmap(matrix=mat)
        slide.shapes.add_picture(io.BytesIO(pix.tobytes("png")), 0, 0, w_emu, h_emu)

    doc.close()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def to_pptx_editable(pdf_bytes: bytes) -> bytes:
    """將文字與圖片提取為獨立圖形，可在 PowerPoint 中直接編輯。"""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    prs = Presentation()

    for page in doc:
        w_emu = _emu(page.rect.width)
        h_emu = _emu(page.rect.height)
        prs.slide_width = w_emu
        prs.slide_height = h_emu
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        page_dict = page.get_text("dict")

        for block in page_dict.get("blocks", []):
            bbox = block["bbox"]
            left   = _emu(bbox[0])
            top    = _emu(bbox[1])
            width  = _emu(bbox[2] - bbox[0])
            height = _emu(bbox[3] - bbox[1])

            if width < 1 or height < 1:
                continue

            if block["type"] == 1:  # 圖片區塊
                try:
                    img_data = block.get("image")
                    if not img_data:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=fitz.Rect(bbox))
                        img_data = pix.tobytes("png")
                    slide.shapes.add_picture(io.BytesIO(img_data), left, top, width, height)
                except Exception:
                    pass
                continue

            if block["type"] != 0:  # 非文字區塊略過
                continue

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = False
            txBox.fill.background()
            txBox.line.fill.background()

            first_para = True
            for line in block.get("lines", []):
                if first_para:
                    para = tf.paragraphs[0]
                    first_para = False
                else:
                    para = tf.add_paragraph()
                para.space_before = Pt(0)
                para.space_after = Pt(0)

                for span in line.get("spans", []):
                    text = span.get("text", "")
                    if not text:
                        continue
                    run = para.add_run()
                    run.text = text
                    font = run.font
                    font.size = Pt(max(span.get("size", 10), 6))

                    color = span.get("color", 0)
                    if isinstance(color, int):
                        font.color.rgb = RGBColor(
                            (color >> 16) & 0xFF,
                            (color >> 8) & 0xFF,
                            color & 0xFF,
                        )

                    flags = span.get("flags", 0)
                    font.bold   = bool(flags & 16)
                    font.italic = bool(flags & 2)

                    raw_name = span.get("font", "")
                    if raw_name:
                        clean = raw_name.split("+")[-1].split(",")[0]
                        try:
                            font.name = clean
                        except Exception:
                            pass

    doc.close()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── conversion：圖片（ZIP）────────────────────────────────────────────────────

def to_images_zip(pdf_bytes: bytes, fmt: str, dpi: int) -> bytes:
    """將每頁轉為圖片，打包成 ZIP 下載。"""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    scale = dpi / 72
    mat = fitz.Matrix(scale, scale)

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            if fmt == "jpg":
                img_bytes = pix.tobytes("jpeg")
                zf.writestr(f"page_{i + 1:03d}.jpg", img_bytes)
            else:
                img_bytes = pix.tobytes("png")
                zf.writestr(f"page_{i + 1:03d}.png", img_bytes)

    doc.close()
    zip_buf.seek(0)
    return zip_buf.getvalue()


# ── UI ────────────────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="PDF 轉檔工具",
    page_icon="📊",
    layout="wide",
)

st.title("📊 PDF 轉檔工具")
st.caption("上傳 PDF，選擇輸出格式後下載，或直接匯入 Canva 進行編輯")

uploaded = st.file_uploader("上傳 PDF 檔案", type=["pdf"])

if not uploaded:
    st.info("👆 請上傳一個 PDF 檔案以開始使用。")
    st.stop()

pdf_bytes = uploaded.read()

try:
    n_pages = pdf_page_count(pdf_bytes)
except Exception as e:
    st.error(f"無法開啟此 PDF：{e}")
    st.stop()

st.success(f"✅ **{uploaded.name}** — 共 {n_pages} 頁")

# ── 頁面預覽 ──────────────────────────────────────────────────────────────────
with st.expander("🔍 頁面預覽", expanded=True):
    previews = render_previews(pdf_bytes)
    cols = st.columns(min(len(previews), 3))
    for i, img in enumerate(previews):
        cols[i % 3].image(img, caption=f"第 {i + 1} 頁", use_container_width=True)
    if n_pages > 6:
        st.caption(f"顯示前 6 頁，共 {n_pages} 頁")

st.divider()

col_convert, col_canva = st.columns(2, gap="large")

# ── 左欄：格式選擇與轉換 ──────────────────────────────────────────────────────
with col_convert:
    st.subheader("⬇️ 選擇輸出格式")

    fmt = st.radio(
        "輸出格式",
        ["pptx", "png", "jpg"],
        format_func=lambda x: {
            "pptx": "📑 PPTX — PowerPoint 投影片",
            "png":  "🖼️ PNG  — 高品質無損圖片（每頁一張）",
            "jpg":  "📷 JPG  — 壓縮圖片，檔案較小（每頁一張）",
        }[x],
        horizontal=False,
    )

    stem = uploaded.name.rsplit(".", 1)[0]

    # PPTX 子選項
    if fmt == "pptx":
        st.markdown("**投影片模式**")
        mode = st.radio(
            "選擇轉換模式",
            ["image", "editable"],
            format_func=lambda x: {
                "image":    "🖼️ 圖片模式 — 版面完全保留，適合直接簡報",
                "editable": "✏️ 可編輯模式 — 文字與圖片皆可在 PowerPoint 中修改",
            }[x],
        )
        st.caption(
            "圖片模式保留每個像素的排版；"
            "可編輯模式將文字提取為獨立文字方塊，可直接點擊修改。"
        )

    # PNG / JPG 子選項
    else:
        st.markdown("**圖片解析度（DPI）**")
        dpi = st.select_slider(
            "解析度",
            options=[72, 96, 150, 200, 300],
            value=150,
            format_func=lambda x: {
                72:  "72 DPI — 螢幕預覽用",
                96:  "96 DPI — 一般品質",
                150: "150 DPI — 建議（平衡品質與檔案大小）",
                200: "200 DPI — 高品質",
                300: "300 DPI — 印刷品質（檔案較大）",
            }[x],
        )
        st.caption(f"共 {n_pages} 張圖片，將打包為一個 ZIP 檔案下載。")

    st.write("")

    if st.button("🔄 開始轉換", type="primary", use_container_width=True):
        try:
            if fmt == "pptx":
                label = "圖片投影片" if mode == "image" else "可編輯投影片"
                with st.spinner(f"正在建立 {label}（共 {n_pages} 頁）…"):
                    data = to_pptx_image(pdf_bytes) if mode == "image" else to_pptx_editable(pdf_bytes)
                st.download_button(
                    label="📥 下載 .pptx",
                    data=data,
                    file_name=f"{stem}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
                if fmt == "pptx" and mode == "editable":
                    st.info("💡 版面為近似還原，複雜排版可能略有偏差。如需版面完整且可編輯，建議使用右側的 Canva 方案。")

            else:
                fmt_label = "PNG" if fmt == "png" else "JPG"
                with st.spinner(f"正在轉換為 {fmt_label}（{dpi} DPI，共 {n_pages} 頁）…"):
                    data = to_images_zip(pdf_bytes, fmt=fmt, dpi=dpi)
                st.download_button(
                    label=f"📥 下載 ZIP（{n_pages} 張 {fmt_label}）",
                    data=data,
                    file_name=f"{stem}_{fmt_label}.zip",
                    mime="application/zip",
                    use_container_width=True,
                )

        except Exception as e:
            st.error(f"轉換失敗：{e}")

# ── 右欄：Canva ───────────────────────────────────────────────────────────────
with col_canva:
    st.subheader("🎨 在 Canva 中編輯")
    st.markdown(
        "Canva 可將 PDF 匯入為**完全可編輯的設計**，"
        "文字、圖片、版面皆完整保留，且每個元素都可以自由修改。"
    )

    with st.container(border=True):
        st.markdown("""
**匯入 Canva 的步驟：**

1. 點擊下方「下載 PDF」
2. 點擊下方「開啟 Canva」
3. 在 Canva 點選 **「建立設計」** → **「匯入檔案」**
4. 選擇剛才下載的 PDF
5. 每一頁都會變成可編輯的投影片 ✅
""")

    c1, c2 = st.columns(2)
    with c1:
        st.download_button(
            label="📥 下載 PDF",
            data=pdf_bytes,
            file_name=uploaded.name,
            mime="application/pdf",
            use_container_width=True,
        )
    with c2:
        st.link_button(
            "🎨 開啟 Canva",
            "https://www.canva.com/",
            use_container_width=True,
        )

    st.markdown("---")
    st.markdown("""
**在 Canva 中你可以：**
- ✏️ 點擊任何文字直接修改內容
- 🖼️ 點擊圖片進行替換、裁切或縮放
- 🎨 自由更換顏色、字體與背景
- 📐 拖曳元素調整位置，版面不會跑掉
- 📤 匯出為 PDF、PPTX 或 PNG
""")
